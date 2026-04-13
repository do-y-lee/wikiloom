"""Office document extractor for .docx and .pptx files."""

from __future__ import annotations

from pathlib import Path

from wikiloom.ingest.extractors.base import BaseExtractor, ExtractedContent
from wikiloom.utils import estimate_tokens

OFFICE_EXTENSIONS = {".docx", ".pptx"}


class OfficeExtractor(BaseExtractor):
    """Extract text from Word (.docx) and PowerPoint (.pptx) files."""

    def can_handle(self, path: Path) -> bool:
        return path.suffix.lower() in OFFICE_EXTENSIONS

    def extract(self, path: Path) -> ExtractedContent:
        ext = path.suffix.lower()
        if ext == ".docx":
            text, method = self._extract_docx(path)
        elif ext == ".pptx":
            text, method = self._extract_pptx(path)
        else:
            raise ValueError(f"OfficeExtractor cannot handle {ext}")

        return ExtractedContent(
            text=text,
            metadata={"filename": path.name},
            source_path=path,
            content_type="office",
            extraction_method=method,
            token_estimate=estimate_tokens(text),
        )

    def _extract_docx(self, path: Path) -> tuple[str, str]:
        import docx  # python-docx

        document = docx.Document(str(path))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs), "python-docx"

    def _extract_pptx(self, path: Path) -> tuple[str, str]:
        from pptx import Presentation  # python-pptx

        prs = Presentation(str(path))
        slide_texts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts = [f"# Slide {idx}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            slide_texts.append("\n".join(parts))
        return "\n\n".join(slide_texts), "python-pptx"
